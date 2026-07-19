import io

import pandas as pd
import streamlit as st

from utils.file_ops import init_module_state

PREFIX = "module5_ppt"


def run():
    st.title("📊 PPT 处理及生成")
    st.caption("创建和编辑 PowerPoint 演示文稿")

    init_module_state(PREFIX, {"slides": []})

    st.subheader("创建 PPT 演示文稿")
    
    presentation_title = st.text_input("演示文稿标题", value="演示报告", key=f"{PREFIX}_title")

    if st.button("添加幻灯片", key=f"{PREFIX}_add_slide"):
        st.session_state[f"{PREFIX}_slides"].append({"title": "", "content": ""})

    slides = st.session_state.get(f"{PREFIX}_slides", [])
    for i, slide in enumerate(slides):
        st.markdown(f"**幻灯片 {i+1}**")
        col1, col2 = st.columns([2, 1])
        with col1:
            slides[i]["title"] = st.text_input(f"标题 {i+1}", value=slide["title"], key=f"{PREFIX}_slide_title_{i}")
            slides[i]["content"] = st.text_area(f"内容 {i+1}", value=slide["content"], key=f"{PREFIX}_slide_content_{i}", height=100)
        with col2:
            if st.button(f"删除幻灯片 {i+1}", key=f"{PREFIX}_delete_slide_{i}"):
                slides.pop(i)
                st.rerun()

    if slides:
        if st.button("生成 PPT", key=f"{PREFIX}_generate", type="primary"):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN

                prs = Presentation()

                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = presentation_title
                subtitle.text = "业财融合多功能系统"

                for i, slide_data in enumerate(slides):
                    bullet_slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    
                    title_shape = slide.shapes.title
                    body_shape = slide.placeholders[1]
                    
                    title_shape.text = slide_data["title"]
                    
                    tf = body_shape.text_frame
                    tf.text = slide_data["content"]
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(14)

                output = io.BytesIO()
                prs.save(output)
                output.seek(0)

                st.success("PPT 生成成功！")
                st.download_button(
                    label="下载 PPT",
                    data=output.getvalue(),
                    file_name=f"{presentation_title}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"{PREFIX}_download",
                )
            except ImportError:
                st.warning("需要安装 python-pptx 库，请运行 pip install python-pptx")
            except Exception as e:
                st.error(f"生成失败：{str(e)}")
    else:
        st.info("请先添加幻灯片")

    st.subheader("PPT 模板")
    st.info("💡 提示：后续可扩展支持上传 PPT 模板并填充数据")